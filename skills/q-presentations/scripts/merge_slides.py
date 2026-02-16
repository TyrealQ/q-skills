"""Merge slide PNGs into PPTX."""
import sys
import glob
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu


def merge_to_pptx(slide_dir: str):
    slide_dir = Path(slide_dir)
    if not slide_dir.exists():
        print(f"Error: Directory not found: {slide_dir}")
        return

    pngs = sorted(glob.glob(str(slide_dir / "*.png")))
    if not pngs:
        print(f"No slide images found in {slide_dir}")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    for png in pngs:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            png, Emu(0), Emu(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
        print(f"Added: {Path(png).name}")

    output = slide_dir / f"{slide_dir.name}.pptx"
    prs.save(str(output))
    print(f"\nCreated: {output}")
    print(f"Total slides: {len(pngs)}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: merge_slides.py <slide_dir>")
        sys.exit(1)
    merge_to_pptx(sys.argv[1])
